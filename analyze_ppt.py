import os
import re
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()


genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')



def read_txt_with_refs(file_path):
    statements = []
    with open(file_path, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            ref_match = re.match(r'(Page|Slide)\s*(\d+):\s*(.*)', line)
            if ref_match:
                statements.append({
                    "reference": f"{ref_match.group(1)} {ref_match.group(2)}",
                    "text": ref_match.group(3)
                })
            else:
                statements.append({"reference": None, "text": line})
    return statements


def read_pptx_with_refs(file_path):
    prs = Presentation(file_path)
    statements = []
    for i, slide in enumerate(prs.slides, start=1):
        text = "\n".join(
            [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        )
        if text.strip():
            statements.append({"reference": f"Slide {i}", "text": text.strip()})
    return statements


def build_faiss_index(statements):
    texts = [s["text"] for s in statements]
    embeddings = embedding_model.encode(texts, convert_to_numpy=True)
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    return index, embeddings


def get_top_related_pairs(statements, embeddings, top_k=5):
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    _, idxs = index.search(embeddings, top_k + 1)

    pairs = set()
    for i, neighbors in enumerate(idxs):
        for j in neighbors[1:]:  # skip itself
            if i < j:  # avoid duplicates
                pairs.add((i, j))
    return list(pairs)


def has_number_mismatch(s1, s2):
    nums1 = re.findall(r'\d+(?:\.\d+)?', s1)
    nums2 = re.findall(r'\d+(?:\.\d+)?', s2)
    if nums1 and nums2 and set(nums1) != set(nums2):
        return True
    return False


def batch_detect_contradictions(pairs, statements):
    model = genai.GenerativeModel("gemini-2.5-flash")
    prompt = """
    You are an AI that detects factual/logical contradictions between pairs of statements.
    For each pair, return item with and when contradiction is True:
    {
        "contradiction": true/false,
        "type": "Numerical" | "Timeline" | "Qualitative" | null,
        "reason": "short explanation",
        "reference_1": "...",
        "reference_2": "..."
    }
    Analyze the following pairs:
    [
    """
    for i, j in pairs:
        prompt += f"""{{
            "s1": "{statements[i]['text']}",
            "ref1": "{statements[i]['reference']}",
            "s2": "{statements[j]['text']}",
            "ref2": "{statements[j]['reference']}"
        }},
        """
    prompt += "]"

    resp = model.generate_content(prompt)
    return resp

def analyze_document(file_path):
  ext = os.path.splitext(file_path)[1].lower()
  if ext == ".pptx":
      statements = read_pptx_with_refs(file_path)
  elif ext == ".txt":
      statements = read_txt_with_refs(file_path)
  else:
      raise ValueError("Unsupported file type. Please use .pptx or .txt")

  index, embeddings = build_faiss_index(statements)
  pair_idxs = get_top_related_pairs(statements, embeddings, top_k=100)

  pairs = [(statements[i], statements[j]) for i, j in pair_idxs]
      # print("ss  ",pairs)
  results = batch_detect_contradictions(pair_idxs,statements)
  return results
import json
if __name__ == "__main__":
  file_path = "NoogatAssignment.pptx"  # or .txt
  result = analyze_document(file_path)
  print(result.text)

